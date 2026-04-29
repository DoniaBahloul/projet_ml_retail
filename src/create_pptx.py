"""
Génère un PPTX de 20 slides :
- Slides 1-7 : images téléchargées depuis les fichiers HTML existants
- Slides 8-20 : générées avec python-pptx (style dark mode + images reports/)
"""
import re
import urllib.request
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Couleurs ──────────────────────────────────────────────────────────────────
BG   = RGBColor(0x0F, 0x0F, 0x1A)
CARD = RGBColor(0x1E, 0x1E, 0x32)
CARD2= RGBColor(0x16, 0x16, 0x28)
ACC  = RGBColor(0x6C, 0x63, 0xFF)
CYA  = RGBColor(0x00, 0xD4, 0xFF)
WHT  = RGBColor(0xFF, 0xFF, 0xFF)
GRY  = RGBColor(0xB0, 0xB0, 0xC0)
GRN  = RGBColor(0x00, 0xE0, 0x96)

W = Inches(13.33)
H = Inches(7.5)
REPORTS = Path("reports")
TMP = Path("tmp_slides")
TMP.mkdir(exist_ok=True)

# ── Utilitaires ───────────────────────────────────────────────────────────────
def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, txt, l, t, w, h, sz=16, bold=False, col=WHT,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = col

def rect(slide, l, t, w, h, col):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = col
    s.line.fill.background()

def header(slide, title, subtitle, num, total=20):
    set_bg(slide)
    rect(slide, 0, 0, Inches(0.07), H, ACC)
    box(slide, f"{num:02d}/{total}", Inches(12.3), Inches(0.1),
        Inches(0.9), Inches(0.4), sz=11, col=GRY, align=PP_ALIGN.RIGHT)
    box(slide, title, Inches(0.25), Inches(0.2),
        Inches(12.8), Inches(0.85), sz=28, bold=True, col=WHT)
    rect(slide, Inches(0.25), Inches(1.05), Inches(3.5), Emu(3500), ACC)
    if subtitle:
        box(slide, subtitle, Inches(0.25), Inches(1.15),
            Inches(12.8), Inches(0.5), sz=15, italic=True, col=CYA)

def bullets(slide, items, top=Inches(1.75)):
    bh = Inches(0.44)
    for i, item in enumerate(items):
        c = CARD if i % 2 == 0 else CARD2
        rect(slide, Inches(0.22), top, Inches(12.9), bh, c)
        box(slide, "▸", Inches(0.32), top+Inches(0.06), Inches(0.3), bh,
            sz=13, col=ACC, bold=True)
        box(slide, item, Inches(0.62), top+Inches(0.04), Inches(12.5), bh,
            sz=12.5, col=WHT)
        top += bh + Inches(0.05)

def add_img(slide, path, l=Inches(0.3), t=Inches(1.7), w=Inches(12.7)):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), l, t, width=w)

# ── Étape 1 : Télécharger les images des slides 1-7 ─────────────────────────
html_files = [f"page_{i}.html" for i in range(1, 8)]
downloaded = []
headers = {"User-Agent": "Mozilla/5.0"}

for i, hf in enumerate(html_files, start=1):
    out_png = TMP / f"slide_{i:02d}.png"
    downloaded.append(out_png)
    if out_png.exists():
        print(f"  [cache] {out_png.name}")
        continue
    try:
        content = Path(hf).read_text(encoding="utf-8")
        url = re.search(r'src="(https://[^"]+\.png)"', content)
        if url:
            url = url.group(1)
            req = urllib.request.Request(url, headers=headers)
            with urllib.request.urlopen(req, timeout=15) as resp:
                out_png.write_bytes(resp.read())
            print(f"  [ok] {out_png.name}")
        else:
            print(f"  [skip] no URL in {hf}")
            downloaded[-1] = None
    except Exception as e:
        print(f"  [err] {hf}: {e}")
        downloaded[-1] = None

# ── Étape 2 : Construire le PPTX ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# --- Slides 1-7 : image plein écran ------------------------------------------
for i, img_path in enumerate(downloaded, start=1):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, RGBColor(0,0,0))
    if img_path and img_path.exists():
        sl.shapes.add_picture(str(img_path), 0, 0, width=W, height=H)
    else:
        box(sl, f"Slide {i} — image non disponible",
            Inches(1), Inches(3), Inches(11), Inches(1),
            sz=24, col=GRY, align=PP_ALIGN.CENTER)
    print(f"  Slide {i} ajoutée")

# --- Slide 8 : Découvertes EDA -----------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Découvertes Clés de l'EDA", "Ce que les graphiques révèlent sur nos données", 8)
bullets(sl, [
    "Déséquilibre de classes : ~18% Churn vs 82% Fidèles → problème de classification !",
    "Corrélations très fortes (|r| > 0.85) entre les features monétaires → multicolinéarité",
    "SupportTickets : valeurs aberrantes allant jusqu'à 999 (erreurs de saisie système)",
    "Satisfaction : valeurs négatives (-1) et extrêmes (99) = codes d'erreur à corriger",
    "Recency et Frequency : fortement corrélés au Churn → features les plus prédictives",
    "Feature 'Newsletter' : identique ('Yes') pour 100% des clients → variance nulle → à supprimer",
])

# --- Slide 9 : Preprocessing Nettoyage ---------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Preprocessing (1/2) — Nettoyage", "src/preprocessing.py", 9)
bullets(sl, [
    "Parsing RegistDate : pd.to_datetime() multi-format → RegYear, RegMonth, RegDay, RegWeekday",
    "Suppression variance nulle : 'Newsletter' supprimé (identique pour tous les clients)",
    "Correction outliers : SupportTickets > Q3+3×IQR → remplacé par la médiane",
    "Correction outliers : Satisfaction < 0 ou > 10 → remplacé par la médiane",
    "Imputation NaN : Age imputé par la médiane (292 valeurs manquantes corrigées)",
    "Suppression multicolinéarité : 5 features supprimées (seuil |r| > 0.85)",
])

# --- Slide 10 : Preprocessing Feature Engineering ---------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Preprocessing (2/2) — Feature Engineering & Encodage",
       "52 features brutes → 74 features propres", 10)
bullets(sl, [
    "IP Engineering : LastLoginIP → IP_IsPrivate (0/1) + IP_FirstOctet (numérique)",
    "RFM Features : MonetaryPerDay, AvgBasketValue, TenureRatio, FrequencyRatio",
    "Encodage Ordinal : AgeCategory, SpendingCat, LoyaltyLevel (ordre logique préservé)",
    "One-Hot Encoding : Gender, Region, FavoriteSeason, PreferredTime… (pd.get_dummies)",
    "Target Encoding : Country → taux de Churn moyen par pays (évite explosion de colonnes)",
    "RÈGLE D'OR — StandardScaler : fit() sur X_train UNIQUEMENT → pas de Data Leakage !",
])

# --- Slide 11 : Split & Scaling ----------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Split Train/Test & Normalisation", "Prévention absolue du Data Leakage", 11)
bullets(sl, [
    "Ordre obligatoire : Séparer (split) AVANT de normaliser (scaler) — jamais l'inverse !",
    "train_test_split(test_size=0.20, stratify=y) → 800 train / 200 test",
    "stratify=y garantit 18% de Churn dans train ET dans test (split équitable)",
    "StandardScaler (Z-Score) : centre à 0, variance=1 → obligatoire pour KMeans et Lasso",
    "Scaler sauvegardé : models/scaler.joblib → réutilisé à la prédiction (cohérence garantie)",
    "Résultat final : X_train (800×74) | X_test (200×74) | y_train / y_test sauvegardés",
])

# --- Slide 12 : KMeans -------------------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Modèle 1 — Clustering KMeans", "Segmentation non supervisée des clients", 12)
add_img(sl, REPORTS/"clustering_elbow.png",  Inches(0.2), Inches(1.6), Inches(6.4))
add_img(sl, REPORTS/"clustering_pca2d.png",  Inches(6.8), Inches(1.6), Inches(6.3))
box(sl, "Méthode Elbow + Score Silhouette → k=2 clusters optimaux  |  Projection ACP 2D des groupes",
    Inches(0.2), Inches(6.9), Inches(12.9), Inches(0.45),
    sz=13, col=CYA, align=PP_ALIGN.CENTER, italic=True)

# --- Slide 13 : Random Forest ------------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Modèle 2 — Classification Churn", "Random Forest + GridSearchCV", 13)
bullets(sl, [
    "Déséquilibre : 18% Churn vs 82% Fidèles → class_weight='balanced'",
    "GridSearchCV 5-fold : n_estimators=[100,200] × max_depth=[10,20] → 20 combinaisons testées",
    "Best params : n_estimators=100, max_depth=10 | F1 CV = 0.84",
], top=Inches(1.75))
add_img(sl, REPORTS/"churn_rf_evaluation.png", Inches(0.2), Inches(3.5), Inches(12.9))

# --- Slide 14 : Feature Importance -------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Évaluation — Feature Importances RF", "Quelles features prédisent le mieux le Churn ?", 14)
add_img(sl, REPORTS/"churn_feature_importance.png", Inches(1.8), Inches(1.55), Inches(9.7))
box(sl, "Recency, Frequency, MonetaryTotal et Satisfaction dominent la prédiction du Churn",
    Inches(0.2), Inches(6.9), Inches(12.9), Inches(0.45),
    sz=13, col=CYA, align=PP_ALIGN.CENTER, italic=True)

# --- Slide 15 : Régression ---------------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Modèle 3 — Régression MonetaryTotal", "Ridge (L2) vs Lasso (L1)", 15)
bullets(sl, [
    "Ridge L2 : Réduit les coefficients sans les annuler → robuste face à la multicolinéarité",
    "Lasso L1 : Annule les coefficients faibles → sélection automatique de features",
    "Validation croisée 5-fold sur X_train | Métriques : MAE, RMSE, R²",
], top=Inches(1.75))
add_img(sl, REPORTS/"regression_monetary_evaluation.png", Inches(0.2), Inches(3.5), Inches(12.9))

# --- Slide 16 : Flask --------------------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Déploiement — Flask Web App", "app/app.py + src/predict.py", 16)
bullets(sl, [
    "Route GET  '/'           → Formulaire de saisie client (16 champs) — index.html",
    "Route POST '/predict'    → Reçoit le formulaire → appelle predict_customer() → result.html",
    "Route GET  '/dashboard'  → Affiche les graphiques d'évaluation depuis reports/*.png",
    "Route POST '/api/predict'→ API REST JSON (intégrable dans n'importe quel système externe)",
    "predict.py : charge scaler + RF + KMeans → normalise → prédit Churn + Cluster + Risque",
    "Logique métier : P>0.75 → 🔴 Critique (intervention immédiate) | P<0.20 → 🟢 Fidèle",
])

# --- Slide 17 : predict.py détail -------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Zoom — Module predict.py", "Le pont entre l'IA et l'utilisateur", 17)
bullets(sl, [
    "load_pipeline() : charge scaler + churn_rf + kmeans + feature_names depuis models/",
    "preprocess_input() : aligne les features du formulaire sur les 74 colonnes attendues",
    "scaler.transform() : applique la même normalisation qu'à l'entraînement (cohérence)",
    "churn_model.predict_proba()[:,1] : retourne la probabilité continue de churn [0, 1]",
    "kmeans.predict() : attribue le cluster le plus proche mathématiquement au nouveau client",
    "Retour dict : churn_proba, risk_level, risk_emoji, risk_message, cluster_label",
])

# --- Slide 18 : Résultats ----------------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Résultats & Métriques Finales", "Évaluation des 3 modèles sur X_test (200 clients)", 18)
bullets(sl, [
    "KMeans       : k=2 clusters | Silhouette Score=0.034 | Distribution : {0:461, 1:339}",
    "Random Forest: Accuracy=95.5% | F1-Score=0.87 | AUC-ROC=0.99",
    "Lasso        : R²=0.81 | RMSE=1995 (Vainqueur face au Ridge)",
    "Pipeline     : 52 features brutes → 74 features → 3 modèles sauvegardés (.joblib)",
    "Durée totale : Génération 5s | Preprocessing 60s | Training 15s (backend Agg)",
])

# --- Slide 19 : Bilan --------------------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Bilan Technique du Projet", "Ce qui a été implémenté de A à Z", 19)
bullets(sl, [
    "6 scripts Python modulaires et documentés (génération → déploiement)",
    "74 features finales issues de 52 features brutes après Feature Engineering complet",
    "3 modèles ML persistés : KMeans + Random Forest + Lasso (via joblib)",
    "1 application Flask déployée localement avec API REST et interface UI premium",
    "1 Notebook Jupyter complet pour l'EDA (26 cellules de code et visualisations)",
    "Respect des bonnes pratiques : pas de data leakage, validation croisée, stratification",
])

# --- Slide 20 : Conclusion ---------------------------------------------------
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
rect(sl, 0, 0, Inches(0.07), H, GRN)
rect(sl, Inches(0.07), Inches(3.7), W-Inches(0.07), Inches(0.05), ACC)
box(sl, "Conclusion & Perspectives",
    Inches(0.3), Inches(0.25), Inches(12.5), Inches(0.9),
    sz=30, bold=True, col=WHT)
rect(sl, Inches(0.25), Inches(1.1), Inches(3.0), Emu(3500), GRN)
box(sl, "Compétences validées",
    Inches(0.3), Inches(1.2), Inches(12.5), Inches(0.45),
    sz=15, col=GRN, italic=True)
bullets(sl, [
    "ETL & Feature Engineering  |  Gestion du déséquilibre de classes (class_weight)",
    "Validation croisée & GridSearchCV  |  Prévention absolue du Data Leakage",
    "Persistance des modèles (joblib)  |  Déploiement Flask + API REST",
], top=Inches(1.75))
box(sl, "Perspectives",
    Inches(0.3), Inches(4.0), Inches(12.5), Inches(0.45),
    sz=15, col=CYA, italic=True)
bullets(sl, [
    "Dockerisation de l'app Flask pour déploiement cloud (AWS / Azure / Heroku)",
    "Connexion à une vraie BDD PostgreSQL et retraining automatique (CI/CD GitHub Actions)",
    "Monitoring de la dérive du modèle (Data Drift Detection) en production",
], top=Inches(4.5))

# ── Sauvegarde ────────────────────────────────────────────────────────────────
out = "presentation_ml_retail_20slides.pptx"
prs.save(out)
print(f"\nPresentation saved: {out}  ({len(prs.slides)} slides)")
